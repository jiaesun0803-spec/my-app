import streamlit as st
import time
import json
import os
import io
import re
import numpy as np
import google.generativeai as genai
import plotly.graph_objects as go
from datetime import datetime

st.set_page_config(page_title="AI 컨설팅 시스템", layout="wide")

# ==========================================
# 유틸리티 함수
# ==========================================
def safe_int(value):
    try:
        clean_val = str(value).replace(',', '').strip()
        if not clean_val: return 0
        return int(float(clean_val))
    except:
        return 0

def format_kr_currency(value):
    try:
        val = safe_int(value)
        if val == 0: return "0원"
        uk = val // 10000
        man = val % 10000
        if uk > 0 and man > 0: return f"{uk}억 {man:,}만원"
        elif uk > 0: return f"{uk}억원"
        else: return f"{man:,}만원"
    except:
        return str(value)

def format_biz_no(raw_no):
    no = str(raw_no).replace("-", "").strip()
    if len(no) == 10: return f"{no[:3]}-{no[3:5]}-{no[5:]}"
    return raw_no

def format_corp_no(raw_no):
    no = str(raw_no).replace("-", "").strip()
    if len(no) == 13: return f"{no[:6]}-{no[6:]}"
    return raw_no

def get_credit_grade(score, type="NICE"):
    score = safe_int(score)
    if type == "NICE":
        if score >= 900: return 1
        elif score >= 870: return 2
        elif score >= 840: return 3
        elif score >= 805: return 4
        elif score >= 750: return 5
        elif score >= 665: return 6
        elif score >= 600: return 7
        elif score >= 515: return 8
        elif score >= 445: return 9
        else: return 10
    else:
        if score >= 942: return 1
        elif score >= 891: return 2
        elif score >= 832: return 3
        elif score >= 768: return 4
        elif score >= 698: return 5
        elif score >= 630: return 6
        elif score >= 530: return 7
        elif score >= 454: return 8
        elif score >= 335: return 9
        else: return 10

# ==========================================
# 업체 DB
# ==========================================
DB_FILE = "company_db.json"

def load_db():
    if os.path.exists(DB_FILE):
        try:
            with open(DB_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except:
            return {}
    return {}

def save_db(db_data):
    with open(DB_FILE, "w", encoding="utf-8") as f:
        json.dump(db_data, f, ensure_ascii=False, indent=4)

# ==========================================
# 후처리: 마침표 줄바꿈 + 하이픈 (HTML용)
# ==========================================
def format_text_to_bullets(text):
    """순수 텍스트를 마침표 기준으로 나눠 - 하이픈 붙인 후 <br> 연결"""
    text = re.sub(r'^[-–•]\s*', '', text.strip())
    sentences = re.split(r'(?<=[.!?])\s+', text.strip())
    result = []
    for s in sentences:
        s = re.sub(r'^[-–•]+\s*', '', s.strip())
        if len(s) > 3:
            result.append(f"- {s}")
    return '<br>'.join(result) if result else text

def postprocess_html(html_text):
    """HTML div 내부 텍스트에 마침표 줄바꿈 + 하이픈 강제 적용"""
    def process_div(m):
        open_tag = m.group(1)
        content  = m.group(2)
        close    = m.group(3)
        if len(re.findall(r'<(?!br)[^>]+>', content)) > 3:
            return m.group(0)
        plain = re.sub(r'<br\s*/?>', ' ', content, flags=re.IGNORECASE)
        plain = re.sub(r'<[^>]+>', '', plain)
        plain = re.sub(r'&[a-z]+;', ' ', plain)
        plain = re.sub(r'\s+', ' ', plain).strip()
        if len(plain) < 10:
            return m.group(0)
        return f"{open_tag}{format_text_to_bullets(plain)}{close}"

    return re.sub(
        r'(<div[^>]*>)((?:(?!<div|</div>).)+)(</div>)',
        process_div, html_text, flags=re.DOTALL
    )

# ==========================================
# 월별 차트 HTML (다운로드용)
# ==========================================
def make_chart_html(monthly_vals, monthly_labels):
    max_v = max(monthly_vals) if monthly_vals else 1
    html = '''<div style="background:#f0f9ff;padding:20px;border-radius:15px;
border:1px solid #BAE6FD;margin:20px 0;page-break-inside:avoid;">
<div style="text-align:center;font-weight:bold;color:#0369A1;
margin-bottom:15px;font-size:17px;">📈 향후 1년간 월별 예상 매출 상승 곡선</div>
<table style="width:100%;height:180px;border-bottom:2px solid #BAE6FD;
border-collapse:collapse;table-layout:fixed;"><tr>'''
    for val in monthly_vals:
        h = int((val / max_v) * 140) + 5
        lbl = format_kr_currency(val).replace('만원','만')
        html += f'''<td style="vertical-align:bottom;padding:0 3px;border:none;">
<div style="font-size:10px;color:#0369A1;margin-bottom:4px;
text-align:center;white-space:nowrap;">{lbl}</div>
<div style="width:80%;height:{h}px;
background:linear-gradient(180deg,#38BDF8,#0EA5E9);
border-radius:4px 4px 0 0;margin:0 auto;"></div></td>'''
    html += '</tr></table><table style="width:100%;border-collapse:collapse;table-layout:fixed;margin-top:5px;"><tr>'
    for lbl in monthly_labels:
        html += f'<td style="text-align:center;font-size:11px;font-weight:bold;color:#0369A1;padding:4px 0;border:none;">{lbl}</td>'
    html += '</tr></table></div>'
    return html

# ==========================================
# HTML 다운로드 생성 (그래프 포함)
# ==========================================
def generate_html_download(c_name, response_text, d, monthly_vals, monthly_labels):
    chart_html = make_chart_html(monthly_vals, monthly_labels)
    body = response_text.replace('[GRAPH_INSERT_POINT]', chart_html)

    return f"""<!DOCTYPE html>
<html lang="ko"><head><meta charset="UTF-8">
<title>{c_name} 기업분석리포트</title>
<style>
*{{box-sizing:border-box;-webkit-print-color-adjust:exact!important;print-color-adjust:exact!important;}}
body{{font-family:'Malgun Gothic','Apple SD Gothic Neo',sans-serif;
padding:40px;line-height:2.0;color:#1e293b;max-width:1050px;
margin:0 auto;font-size:14px;background:#fff;}}
h1{{color:#0369A1;text-align:center;font-size:24px;margin-bottom:6px;}}
.subtitle{{text-align:center;color:#64748b;font-size:13px;margin-bottom:28px;}}
h2,h2.section-title{{color:#0369A1!important;font-size:16px;
border-bottom:2px solid #0369A1!important;padding-bottom:7px;margin-top:28px;}}
table{{width:100%;border-collapse:collapse;margin-bottom:12px;font-size:13px;}}
th,td{{padding:10px 12px;border:1px solid #e2e8f0;vertical-align:top;}}
th{{background:#f0f9ff;font-weight:bold;color:#0369A1;}}
.print-btn{{display:block;width:100%;padding:13px;
background:linear-gradient(135deg,#0EA5E9,#10B981);
color:white;font-size:16px;font-weight:bold;border:none;
border-radius:10px;cursor:pointer;margin-bottom:26px;text-align:center;}}
@media print{{
.print-btn{{display:none;}}
@page{{size:A4;margin:14mm;}}
body{{padding:0!important;font-size:12px!important;}}
h2,h2.section-title{{page-break-before:always;margin-top:0!important;}}
h2.section-title:first-of-type{{page-break-before:avoid;}}}}
</style></head><body>
<button class="print-btn" onclick="window.print()">
🖨️ 클릭하여 PDF로 저장 (Ctrl+P → PDF 선택)
</button>
<h1>📋 AI 기업분석 결과보고서</h1>
<div class="subtitle">
{c_name} &nbsp;|&nbsp; 대표자: {d.get('in_rep_name','미입력')} &nbsp;|&nbsp;
업종: {d.get('in_industry','미입력')} &nbsp;|&nbsp;
작성일: {datetime.now().strftime('%Y년 %m월 %d일')}
</div>
<hr style="border:1.5px solid #0369A1;margin-bottom:26px;">
{body}
</body></html>""".encode('utf-8')

# ==========================================
# PPT 생성 - 하늘+민트 클린 모던
# JSON 데이터 기반으로 100% 내용 반영
# ==========================================
def generate_pptx(c_name, ppt_data, d, monthly_vals, monthly_labels):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        # ── 하늘 + 민트 컬러 팔레트 ──
        SKY    = RGBColor(0x03, 0x69, 0xA1)   # 메인 하늘색
        SKY_L  = RGBColor(0x0E, 0xA5, 0xE9)   # 밝은 하늘
        SKY_XL = RGBColor(0xE0, 0xF2, 0xFE)   # 매우 연한 하늘
        SKY_M  = RGBColor(0xBA, 0xE6, 0xFD)   # 중간 하늘
        MINT   = RGBColor(0x10, 0xB9, 0x81)   # 메인 민트
        MINT_L = RGBColor(0x6E, 0xE7, 0xB7)   # 밝은 민트
        MINT_X = RGBColor(0xD1, 0xFA, 0xE5)   # 매우 연한 민트
        WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
        DARK   = RGBColor(0x1E, 0x29, 0x3B)   # 슬레이트 다크
        GREY   = RGBColor(0x64, 0x74, 0x8B)   # 슬레이트 그레이
        LGREY  = RGBColor(0xF1, 0xF5, 0xF9)   # 연한 회색 배경
        RED    = RGBColor(0xEF, 0x44, 0x44)
        AMBER  = RGBColor(0xF5, 0x9E, 0x0B)

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        def blank():
            return prs.slides.add_slide(prs.slide_layouts[6])

        def set_bg(slide, rgb):
            f = slide.background.fill
            f.solid()
            f.fore_color.rgb = rgb

        def rect(slide, l, t, w, h, fill, line=None, lw=0.5):
            s = slide.shapes.add_shape(
                1, Inches(l), Inches(t), Inches(w), Inches(h))
            s.fill.solid()
            s.fill.fore_color.rgb = fill
            if line:
                s.line.color.rgb = line
                s.line.width = Pt(lw)
            else:
                s.line.fill.background()
            return s

        def txt(slide, text, l, t, w, h,
                size=11, bold=False, color=None,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
            if color is None: color = DARK
            tb = slide.shapes.add_textbox(
                Inches(l), Inches(t), Inches(w), Inches(h))
            tf = tb.text_frame
            tf.word_wrap = wrap
            lines = [ln.strip() for ln in str(text).split('\n') if ln.strip()]
            if not lines: lines = ['']
            first = True
            for line in lines:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.alignment = align
                run = p.add_run()
                run.text = line
                run.font.size  = Pt(size)
                run.font.bold  = bold
                run.font.italic = italic
                run.font.color.rgb = color
            return tb

        def bullet_txt(slide, items, l, t, w, h, size=10.5, color=None):
            """리스트 항목을 - 하이픈으로 출력"""
            if color is None: color = DARK
            tb = slide.shapes.add_textbox(
                Inches(l), Inches(t), Inches(w), Inches(h))
            tf = tb.text_frame
            tf.word_wrap = True
            first = True
            for item in items:
                item = re.sub(r'^[-–•]+\s*', '', str(item).strip())
                if not item: continue
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = f"- {item}"
                run.font.size = Pt(size)
                run.font.color.rgb = color
            return tb

        def header_bar(slide, title, subtitle=""):
            """상단 헤더 바"""
            rect(slide, 0, 0, 13.33, 1.25, SKY)
            rect(slide, 0, 1.25, 13.33, 0.06, MINT)
            txt(slide, title, 0.35, 0.15, 10, 0.75,
                size=22, bold=True, color=WHITE)
            if subtitle:
                txt(slide, subtitle, 0.35, 0.82, 10, 0.38,
                    size=11, color=MINT_L, italic=True)

        def card(slide, l, t, w, h, title="", title_color=None,
                 bg=None, accent=None):
            """카드 박스"""
            if bg is None: bg = SKY_XL
            if title_color is None: title_color = SKY
            rect(slide, l, t, w, h, bg, SKY_M, 0.3)
            if accent:
                rect(slide, l, t, 0.08, h, accent)
            if title:
                txt(slide, title, l+0.18, t+0.1, w-0.25, 0.38,
                    size=11, bold=True, color=title_color)
            return l+0.18, t+0.55

        # ============================================================
        # 슬라이드 1: 표지
        # ============================================================
        s = blank()
        set_bg(s, LGREY)
        # 상단 대형 헤더
        rect(s, 0, 0, 13.33, 4.2, SKY)
        rect(s, 0, 4.2, 13.33, 0.08, MINT)
        # 장식 원형
        from pptx.util import Inches as I
        circ = s.shapes.add_shape(9, I(10.5), I(0.3), I(2.8), I(2.8))
        circ.fill.solid(); circ.fill.fore_color.rgb = SKY_L
        circ.line.fill.background()
        circ2 = s.shapes.add_shape(9, I(11.2), I(-0.3), I(2.0), I(2.0))
        circ2.fill.solid(); circ2.fill.fore_color.rgb = MINT
        circ2.line.fill.background()

        txt(s, "AI 기업분석 결과보고서",
            0.6, 0.7, 10, 1.2, size=36, bold=True, color=WHITE)
        txt(s, c_name, 0.6, 2.0, 10, 0.9, size=28, bold=True,
            color=MINT_L)
        txt(s, f"담당: {d.get('in_rep_name','미입력')}  |  "
            f"업종: {d.get('in_industry','미입력')}",
            0.6, 3.0, 10, 0.5, size=13, color=SKY_M)

        # 하단 정보 카드들
        info_items = [
            ("📅 작성일", datetime.now().strftime('%Y.%m.%d')),
            ("🏢 사업자번호", format_biz_no(d.get('in_raw_biz_no','미입력'))),
            ("💰 필요자금", format_kr_currency(d.get('in_req_amount',0))),
            ("📍 소재지", d.get('in_biz_addr','미입력')[:20]+"..." if len(d.get('in_biz_addr','')) > 20 else d.get('in_biz_addr','미입력')),
        ]
        for i, (label, val) in enumerate(info_items):
            x = 0.3 + i*3.28
            rect(s, x, 4.55, 3.1, 1.4, WHITE, SKY_M, 0.3)
            rect(s, x, 4.55, 3.1, 0.38, SKY_XL)
            txt(s, label, x+0.15, 4.6, 2.8, 0.32,
                size=10, bold=True, color=SKY)
            txt(s, str(val), x+0.15, 5.05, 2.8, 0.75,
                size=11, color=DARK)

        txt(s, "본 보고서는 AI 컨설팅 시스템에 의해 자동 생성되었습니다.",
            0.3, 6.85, 12.7, 0.4, size=10, color=GREY,
            align=PP_ALIGN.CENTER, italic=True)

        # ============================================================
        # 슬라이드 2: 기업현황
        # ============================================================
        s = blank()
        set_bg(s, LGREY)
        header_bar(s, "1. 기업현황분석", "Company Overview")

        corp_items = ppt_data.get('기업현황', [])
        infos = [
            ("기업명",     d.get('in_company_name','미입력')),
            ("대표자명",   d.get('in_rep_name','미입력')),
            ("업종",       d.get('in_industry','미입력')),
            ("사업개시일", d.get('in_start_date','미입력')),
            ("사업자번호", format_biz_no(d.get('in_raw_biz_no','미입력'))),
            ("사업장 주소",d.get('in_biz_addr','미입력')),
        ]
        for i, (label, val) in enumerate(infos):
            col = i % 2
            row = i // 2
            x = 0.25 + col*6.55
            y = 1.45 + row*1.55
            rect(s, x, y, 6.35, 1.38, WHITE, SKY_M, 0.3)
            rect(s, x, y, 6.35, 0.42, SKY_XL)
            txt(s, label, x+0.15, y+0.08, 3, 0.3,
                size=10, bold=True, color=SKY)
            txt(s, str(val), x+0.15, y+0.52, 6.0, 0.72,
                size=11, color=DARK)

        # 분석 내용
        if corp_items:
            rect(s, 0.25, 6.1, 12.85, 1.15, MINT_X, MINT_L, 0.3)
            rect(s, 0.25, 6.1, 0.08, 1.15, MINT)
            txt(s, corp_items[0] if corp_items else '',
                0.45, 6.18, 12.3, 0.95, size=10, color=DARK)

        # ============================================================
        # 슬라이드 3: SWOT
        # ============================================================
        s = blank()
        set_bg(s, LGREY)
        header_bar(s, "2. SWOT 분석", "Strategic Analysis")

        swot = ppt_data.get('swot', {})
        swot_cfg = [
            ('strengths',    'S  강점 (Strengths)',   SKY_XL, SKY,   0.25, 1.45),
            ('weaknesses',   'W  약점 (Weaknesses)',  RGBColor(0xFF,0xED,0xED), RED, 6.8, 1.45),
            ('opportunities','O  기회 (Opportunities)',MINT_X, MINT,  0.25, 4.28),
            ('threats',      'T  위협 (Threats)',     RGBColor(0xFF,0xF7,0xED), AMBER, 6.8, 4.28),
        ]
        for key, label, bgc, tc, x, y in swot_cfg:
            items = swot.get(key, [])
            rect(s, x, y, 6.3, 2.72, bgc, SKY_M, 0.3)
            rect(s, x, y, 0.08, 2.72, tc)
            txt(s, label, x+0.2, y+0.1, 5.8, 0.4,
                size=12, bold=True, color=tc)
            bullet_txt(s, items[:4], x+0.2, y+0.58, 5.9, 2.0,
                       size=10, color=DARK)

        # ============================================================
        # 슬라이드 4: 시장현황
        # ============================================================
        s = blank()
        set_bg(s, LGREY)
        header_bar(s, "3. 시장현황 및 경쟁력", "Market Analysis")

        mkt_items = ppt_data.get('시장현황', [])
        rect(s, 0.25, 1.45, 12.85, 5.8, WHITE, SKY_M, 0.3)
        rect(s, 0.25, 1.45, 0.08, 5.8, SKY)
        txt(s, "📊 시장 현황 분석", 0.48, 1.55, 5, 0.4,
            size=13, bold=True, color=SKY)
        bullet_txt(s, mkt_items[:10], 0.48, 2.05, 12.5, 5.0,
                   size=11, color=DARK)

        # ============================================================
        # 슬라이드 5: 핵심경쟁력
        # ============================================================
        s = blank()
        set_bg(s, LGREY)
        header_bar(s, "4. 핵심경쟁력분석", "Core Competency")

        comp_items = ppt_data.get('핵심경쟁력', [{'제목':'','내용':[]}]*3)
        for i, pt in enumerate(comp_items[:3]):
            y = 1.45 + i*1.95
            rect(s, 0.25, y, 12.85, 1.78, WHITE, SKY_M, 0.3)
            rect(s, 0.25, y, 0.08, 1.78, MINT)
            rect(s, 0.25, y, 12.85, 0.48, MINT_X)
            title_txt = pt.get('제목', f'포인트 {i+1}') if isinstance(pt, dict) else f'포인트 {i+1}'
            txt(s, f"✦  {title_txt}", 0.48, y+0.08, 12, 0.36,
                size=12, bold=True, color=MINT)
            content = pt.get('내용', []) if isinstance(pt, dict) else []
            bullet_txt(s, content[:2], 0.48, y+0.58, 12.2, 1.05,
                       size=10.5, color=DARK)

        # ============================================================
        # 슬라이드 6: 자금사용계획
        # ============================================================
        s = blank()
        set_bg(s, LGREY)
        req_fund  = format_kr_currency(d.get('in_req_amount', 0))
        fund_type = d.get('in_fund_type', '운전자금')
        header_bar(s, f"5. 자금 사용계획  ·  총 {req_fund} ({fund_type})",
                   "Fund Usage Plan")

        fund_items = ppt_data.get('자금계획', [{'항목':'','금액':'','내용':[]}]*2)
        for i, fi in enumerate(fund_items[:2]):
            y = 1.45 + i*2.75
            rect(s, 0.25, y, 12.85, 2.55, WHITE, SKY_M, 0.3)
            rect(s, 0.25, y, 12.85, 0.5, SKY_XL)
            rect(s, 0.25, y, 0.08, 2.55, SKY)
            item_nm  = fi.get('항목','') if isinstance(fi, dict) else ''
            item_amt = fi.get('금액','') if isinstance(fi, dict) else ''
            txt(s, f"{item_nm}   /   {item_amt}", 0.48, y+0.1, 11, 0.36,
                size=12, bold=True, color=SKY)
            content = fi.get('내용', []) if isinstance(fi, dict) else []
            bullet_txt(s, content[:3], 0.48, y+0.65, 12.2, 1.75,
                       size=10.5, color=DARK)

        # ============================================================
        # 슬라이드 7: 매출전망 (4단계)
        # ============================================================
        s = blank()
        set_bg(s, LGREY)
        header_bar(s, "6. 매출 1년 전망", "Sales Forecast")

        stages = ppt_data.get('매출전망', [
            {'단계':'1단계 (도입)','목표':'','내용':[]},
            {'단계':'2단계 (성장)','목표':'','내용':[]},
            {'단계':'3단계 (확장)','목표':'','내용':[]},
            {'단계':'4단계 (안착)','목표':'','내용':[]},
        ])
        step_colors = [
            (SKY_XL, SKY),
            (RGBColor(0xD0,0xF0,0xFD), SKY_L),
            (MINT_X, MINT),
            (RGBColor(0xA7,0xF3,0xD0), RGBColor(0x05,0x96,0x69)),
        ]
        for i, (stage, (bgc, tc)) in enumerate(zip(stages[:4], step_colors)):
            y = 1.45 + i*1.48
            rect(s, 0.25, y, 12.85, 1.35, bgc, SKY_M, 0.3)
            rect(s, 0.25, y, 0.08, 1.35, tc)
            label = stage.get('단계', f'{i+1}단계') if isinstance(stage, dict) else f'{i+1}단계'
            goal  = stage.get('목표', '') if isinstance(stage, dict) else ''
            content = stage.get('내용', []) if isinstance(stage, dict) else []
            txt(s, label, 0.48, y+0.08, 3, 0.38,
                size=12, bold=True, color=tc)
            if goal:
                txt(s, f"🎯 목표: {goal}", 9.5, y+0.08, 3.4, 0.38,
                    size=11, bold=True, color=RED, align=PP_ALIGN.RIGHT)
            if content:
                txt(s, content[0] if content else '',
                    0.48, y+0.58, 12.2, 0.65, size=10, color=DARK)

        # ============================================================
        # 슬라이드 8: 매출 바 차트
        # ============================================================
        s = blank()
        set_bg(s, LGREY)
        header_bar(s, "6. 매출 상승 곡선", "Monthly Sales Chart")

        if monthly_vals:
            max_v = max(monthly_vals) if monthly_vals else 1
            bar_area_h = 4.8
            y_floor = 6.8
            for i, (val, lbl) in enumerate(zip(monthly_vals, monthly_labels)):
                bar_h = (val / max_v) * bar_area_h
                x = 0.4 + i * 1.04
                y_bar = y_floor - bar_h

                # 그라데이션 효과 (두 레이어)
                rect(s, x, y_bar, 0.75, bar_h,
                     SKY_L if i < 6 else MINT)
                rect(s, x, y_bar, 0.18, bar_h, SKY)

                # 값 레이블
                val_str = format_kr_currency(val).replace('만원','만')
                txt(s, val_str, x-0.1, y_bar-0.32, 1.0, 0.3,
                    size=8, color=SKY, align=PP_ALIGN.CENTER)
                # 월 레이블
                txt(s, lbl, x, y_floor+0.05, 0.75, 0.28,
                    size=9, bold=True, color=DARK, align=PP_ALIGN.CENTER)

            # Y축 기준선
            from pptx.util import Pt as Pt2
            line = s.shapes.add_shape(
                1, Inches(0.25), Inches(1.9), Inches(0.02), Inches(5.0))
            line.fill.solid(); line.fill.fore_color.rgb = SKY_M
            line.line.fill.background()

        # ============================================================
        # 슬라이드 9: 성장비전
        # ============================================================
        s = blank()
        set_bg(s, LGREY)
        header_bar(s, "7. 성장비전 및 AI 코멘트", "Growth Vision")

        vision = ppt_data.get('성장비전', {
            '단기': [], '중기': [], '장기': []})
        v_cfg = [
            ('단기', '🌱 단기 비전 (1년)', MINT_X, MINT),
            ('중기', '🚀 중기 비전 (3년)', SKY_XL, SKY),
            ('장기', '👑 장기 비전 (5년)', RGBColor(0xF5,0xF3,0xFF), RGBColor(0x7C,0x3A,0xED)),
        ]
        for i, (key, label, bgc, tc) in enumerate(v_cfg):
            y = 1.45 + i*1.92
            items = vision.get(key, [])
            rect(s, 0.25, y, 12.85, 1.75, bgc, SKY_M, 0.3)
            rect(s, 0.25, y, 0.08, 1.75, tc)
            rect(s, 0.25, y, 12.85, 0.46, RGBColor(
                min(bgc.rgb >> 16, 220),
                min((bgc.rgb >> 8) & 0xFF, 220),
                min(bgc.rgb & 0xFF, 220)
            ) if hasattr(bgc, 'rgb') else bgc)
            txt(s, label, 0.48, y+0.08, 5, 0.36,
                size=12, bold=True, color=tc)
            bullet_txt(s, items[:2], 0.48, y+0.58, 12.2, 1.05,
                       size=10.5, color=DARK)

        # 조언 박스
        advice = ppt_data.get('조언', [])
        rect(s, 0.25, 7.25, 12.85, 0.0, LGREY)  # 여백
        # (성장비전 슬라이드에 조언 추가하기엔 공간 부족 → 별도 슬라이드)

        # ============================================================
        # 슬라이드 10: 인증·특허 조언
        # ============================================================
        s = blank()
        set_bg(s, LGREY)
        header_bar(s, "💡 인증 및 특허 확보 조언", "Certification & IP Strategy")

        advice = ppt_data.get('조언', [])
        rect(s, 0.25, 1.45, 12.85, 5.8, WHITE, SKY_M, 0.3)
        rect(s, 0.25, 1.45, 0.08, 5.8, MINT)
        rect(s, 0.25, 1.45, 12.85, 0.52, MINT_X)
        txt(s, "필수 인증 및 지식재산권 전략", 0.48, 1.55, 8, 0.38,
            size=13, bold=True, color=MINT)
        bullet_txt(s, advice[:12], 0.48, 2.1, 12.3, 5.0,
                   size=11, color=DARK)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.getvalue()

    except Exception as e:
        st.error(f"PPT 생성 오류: {e}")
        import traceback
        st.error(traceback.format_exc())
        return None

# ==========================================
# AI 리포트 + PPT JSON 동시 생성
# ==========================================
def generate_report_and_ppt_data(model, c_name, d):
    c_ind    = d.get('in_industry','미입력')
    rep_name = d.get('in_rep_name','미입력')
    biz_no   = format_biz_no(d.get('in_raw_biz_no','미입력'))
    corp_no  = format_corp_no(d.get('in_raw_corp_no',''))
    corp_text = f" (법인: {corp_no})" if corp_no else ""
    address  = d.get('in_biz_addr','미입력')
    if d.get('in_has_additional_biz') == '유' and d.get('in_additional_biz_addr','').strip():
        address += f" / 추가: {d.get('in_additional_biz_addr')}"
    fund_type = d.get('in_fund_type','운전자금')
    req_fund  = format_kr_currency(d.get('in_req_amount',0))
    item      = d.get('in_item_desc','미입력')

    # ── HTML 리포트 프롬프트 ──
    html_prompt = f"""당신은 20년 경력의 중소기업 경영컨설턴트입니다.
규칙:
1. 문체: '~있음','~예상됨','~확인됨' 간결체. '~습니다' 절대 금지.
2. 마크다운(##,**) 절대 금지. <br> 태그 절대 금지.
3. 각 항목 3~4문장 이상 상세하게 작성.

[기업 정보]
기업명:{c_name} / 대표자:{rep_name} / 업종:{c_ind} / 아이템:{item} / 신청자금:{req_fund}

[출력 - 아래 HTML 구조 그대로]

<h2 class="section-title" style="color:#0369A1;border-bottom:2px solid #0369A1;padding-bottom:8px;margin-top:30px;">1. 기업현황분석</h2>
<table style="width:100%;border-collapse:collapse;font-size:1.05em;background:#f0f9ff;margin-bottom:15px;">
<tr><td style="padding:13px;border-bottom:1px solid #BAE6FD;width:15%;"><b>기업명</b></td><td style="padding:13px;border-bottom:1px solid #BAE6FD;width:35%;">{c_name}</td><td style="padding:13px;border-bottom:1px solid #BAE6FD;width:15%;"><b>대표자명</b></td><td style="padding:13px;border-bottom:1px solid #BAE6FD;width:35%;">{rep_name}</td></tr>
<tr><td style="padding:13px;border-bottom:1px solid #BAE6FD;"><b>업종</b></td><td style="padding:13px;border-bottom:1px solid #BAE6FD;">{c_ind}</td><td style="padding:13px;border-bottom:1px solid #BAE6FD;"><b>사업자번호</b></td><td style="padding:13px;border-bottom:1px solid #BAE6FD;">{biz_no}{corp_text}</td></tr>
<tr><td style="padding:13px;"><b>사업장 주소</b></td><td colspan="3" style="padding:13px;">{address}</td></tr>
</table>
<div style="background:#EFF6FF;padding:15px;border-radius:10px;border-left:4px solid #0EA5E9;margin-bottom:15px;line-height:2.0;">(기업 잠재력 3~4문장)</div>

<h2 class="section-title" style="color:#0369A1;border-bottom:2px solid #0369A1;padding-bottom:8px;margin-top:30px;">2. SWOT 분석</h2>
<table style="width:100%;table-layout:fixed;border-collapse:separate;border-spacing:13px;margin-bottom:15px;">
<tr>
<td style="background:#EFF6FF;padding:18px;border-radius:13px;border-left:4px solid #0369A1;vertical-align:top;line-height:2.0;"><b style="color:#0369A1;">S (강점)</b><div style="margin-top:10px;">(강점 3~4문장)</div></td>
<td style="background:#FEF2F2;padding:18px;border-radius:13px;border-left:4px solid #EF4444;vertical-align:top;line-height:2.0;"><b style="color:#EF4444;">W (약점)</b><div style="margin-top:10px;">(약점 3~4문장)</div></td>
</tr>
<tr>
<td style="background:#ECFDF5;padding:18px;border-radius:13px;border-left:4px solid #10B981;vertical-align:top;line-height:2.0;"><b style="color:#10B981;">O (기회)</b><div style="margin-top:10px;">(기회 3~4문장)</div></td>
<td style="background:#FFFBEB;padding:18px;border-radius:13px;border-left:4px solid #F59E0B;vertical-align:top;line-height:2.0;"><b style="color:#F59E0B;">T (위협)</b><div style="margin-top:10px;">(위협 3~4문장)</div></td>
</tr>
</table>

<h2 class="section-title" style="color:#0369A1;border-bottom:2px solid #0369A1;padding-bottom:8px;margin-top:30px;">3. 시장현황 및 경쟁력 비교</h2>
<div style="background:#EFF6FF;padding:18px;border-radius:13px;border-left:4px solid #0369A1;margin-bottom:13px;line-height:2.0;"><b>📊 시장 현황 분석</b><div style="margin-top:10px;">(시장 트렌드 3~4문장)</div></div>
<div style="padding:13px;background:#fff;border-radius:13px;border:1px solid #BAE6FD;">
<b>⚔️ 주요 경쟁사 비교</b>
<table style="width:100%;border-collapse:collapse;text-align:center;font-size:0.93em;margin-top:10px;">
<tr style="background:#EFF6FF;"><th style="padding:11px;border:1px solid #BAE6FD;">비교 항목</th><th style="padding:11px;border:1px solid #BAE6FD;">{c_name} (자사)</th><th style="padding:11px;border:1px solid #BAE6FD;">경쟁사 A</th><th style="padding:11px;border:1px solid #BAE6FD;">경쟁사 B</th></tr>
<tr><td style="padding:11px;border:1px solid #BAE6FD;font-weight:bold;">핵심 타겟</td><td style="padding:11px;border:1px solid #BAE6FD;">(자사)</td><td style="padding:11px;border:1px solid #BAE6FD;">(A)</td><td style="padding:11px;border:1px solid #BAE6FD;">(B)</td></tr>
<tr><td style="padding:11px;border:1px solid #BAE6FD;font-weight:bold;">차별화 요소</td><td style="padding:11px;border:1px solid #BAE6FD;">(자사)</td><td style="padding:11px;border:1px solid #BAE6FD;">(A)</td><td style="padding:11px;border:1px solid #BAE6FD;">(B)</td></tr>
<tr><td style="padding:11px;border:1px solid #BAE6FD;font-weight:bold;">예상 점유율</td><td style="padding:11px;border:1px solid #BAE6FD;">(자사)</td><td style="padding:11px;border:1px solid #BAE6FD;">(A)</td><td style="padding:11px;border:1px solid #BAE6FD;">(B)</td></tr>
</table>
</div>

<h2 class="section-title" style="color:#0369A1;border-bottom:2px solid #0369A1;padding-bottom:8px;margin-top:30px;">4. 핵심경쟁력분석</h2>
<div style="display:flex;flex-direction:column;gap:12px;margin-bottom:13px;">
<div style="border-left:4px solid #10B981;border-radius:10px;background:#ECFDF5;padding:17px;line-height:2.0;"><div style="font-weight:bold;color:#10B981;margin-bottom:7px;">포인트 1: (키워드)</div><div>(분석 3~4문장)</div></div>
<div style="border-left:4px solid #10B981;border-radius:10px;background:#ECFDF5;padding:17px;line-height:2.0;"><div style="font-weight:bold;color:#10B981;margin-bottom:7px;">포인트 2: (키워드)</div><div>(분석 3~4문장)</div></div>
<div style="border-left:4px solid #10B981;border-radius:10px;background:#ECFDF5;padding:17px;line-height:2.0;"><div style="font-weight:bold;color:#10B981;margin-bottom:7px;">포인트 3: (키워드)</div><div>(분석 3~4문장)</div></div>
</div>

<h2 class="section-title" style="color:#0369A1;border-bottom:2px solid #0369A1;padding-bottom:8px;margin-top:30px;">5. 자금 사용계획 (총 신청자금: {req_fund})</h2>
<table style="width:100%;border-collapse:collapse;text-align:left;margin-bottom:13px;">
<tr style="background:#EFF6FF;"><th style="padding:13px;border:1px solid #BAE6FD;width:20%;">구분 ({fund_type})</th><th style="padding:13px;border:1px solid #BAE6FD;width:60%;">상세 사용계획</th><th style="padding:13px;border:1px solid #BAE6FD;width:20%;">사용예정금액</th></tr>
<tr><td style="padding:13px;border:1px solid #BAE6FD;font-weight:bold;">(세부항목 1)</td><td style="padding:13px;border:1px solid #BAE6FD;line-height:2.0;">(사용처 2~3문장)</td><td style="padding:13px;border:1px solid #BAE6FD;font-weight:bold;color:#0369A1;">(금액)</td></tr>
<tr><td style="padding:13px;border:1px solid #BAE6FD;font-weight:bold;">(세부항목 2)</td><td style="padding:13px;border:1px solid #BAE6FD;line-height:2.0;">(사용처 2~3문장)</td><td style="padding:13px;border:1px solid #BAE6FD;font-weight:bold;color:#0369A1;">(금액)</td></tr>
</table>

<h2 class="section-title" style="color:#0369A1;border-bottom:2px solid #0369A1;padding-bottom:8px;margin-top:30px;">6. 매출 1년 전망</h2>
<div style="display:flex;flex-direction:column;gap:11px;margin-bottom:13px;">
<div style="background:#EFF6FF;padding:17px;border-radius:13px;border-left:4px solid #0369A1;line-height:2.0;"><div style="font-weight:bold;color:#0369A1;margin-bottom:7px;">1단계 (도입)</div><div>(전략 2~3문장)</div><div style="color:#EF4444;font-weight:bold;margin-top:7px;">목표: OOO만원</div></div>
<div style="background:#EFF6FF;padding:17px;border-radius:13px;border-left:4px solid #0369A1;line-height:2.0;"><div style="font-weight:bold;color:#0369A1;margin-bottom:7px;">2단계 (성장)</div><div>(전략 2~3문장)</div><div style="color:#EF4444;font-weight:bold;margin-top:7px;">목표: OOO만원</div></div>
<div style="background:#EFF6FF;padding:17px;border-radius:13px;border-left:4px solid #0369A1;line-height:2.0;"><div style="font-weight:bold;color:#0369A1;margin-bottom:7px;">3단계 (확장)</div><div>(전략 2~3문장)</div><div style="color:#EF4444;font-weight:bold;margin-top:7px;">목표: OOO만원</div></div>
<div style="background:#EFF6FF;padding:17px;border-radius:13px;border-left:4px solid #0369A1;line-height:2.0;"><div style="font-weight:bold;color:#0369A1;margin-bottom:7px;">4단계 (안착)</div><div>(전략 2~3문장)</div><div style="color:#EF4444;font-weight:bold;margin-top:7px;">최종목표: OOO만원</div></div>
</div>

[GRAPH_INSERT_POINT]

<h2 class="section-title" style="color:#0369A1;border-bottom:2px solid #0369A1;padding-bottom:8px;margin-top:30px;">7. 성장비전 및 AI 컨설턴트 코멘트</h2>
<div style="display:flex;flex-direction:column;gap:11px;margin-bottom:18px;">
<div style="background:#ECFDF5;padding:17px;border-radius:13px;border-left:4px solid #10B981;line-height:2.0;"><b style="color:#10B981;">🌱 단기 비전</b><div style="margin-top:9px;">(단기 비전 2~3문장)</div></div>
<div style="background:#EFF6FF;padding:17px;border-radius:13px;border-left:4px solid #0369A1;line-height:2.0;"><b style="color:#0369A1;">🚀 중기 비전</b><div style="margin-top:9px;">(중기 비전 2~3문장)</div></div>
<div style="background:#F5F3FF;padding:17px;border-radius:13px;border-left:4px solid #7C3AED;line-height:2.0;"><b style="color:#7C3AED;">👑 장기 비전</b><div style="margin-top:9px;">(장기 비전 2~3문장)</div></div>
</div>
<div style="background:#ECFDF5;border-left:4px solid #10B981;padding:22px;border-radius:13px;line-height:2.0;">
<b style="color:#10B981;">💡 필수 인증 및 특허 확보 조언</b>
<div style="margin-top:9px;">(인증 조언 2~3문장)</div>
<div style="margin-top:7px;">(지식재산권 전략 2~3문장)</div>
</div>
"""

    # ── PPT JSON 프롬프트 ──
    json_prompt = f"""당신은 중소기업 경영컨설턴트입니다.
아래 기업 정보를 바탕으로 PPT 슬라이드용 JSON 데이터를 생성하세요.
반드시 유효한 JSON만 출력하고, 다른 텍스트는 절대 포함하지 마세요.
문체: '~있음','~예상됨' 간결체. 각 항목은 완전한 문장으로 작성.

기업명:{c_name} / 업종:{c_ind} / 아이템:{item} / 신청자금:{req_fund}

출력 형식 (이 구조 그대로):
{{
  "기업현황": [
    "기업 잠재력 분석 문장1.",
    "기업 잠재력 분석 문장2.",
    "기업 잠재력 분석 문장3."
  ],
  "swot": {{
    "strengths":    ["강점 문장1.", "강점 문장2.", "강점 문장3.", "강점 문장4."],
    "weaknesses":   ["약점 문장1.", "약점 문장2.", "약점 문장3.", "약점 문장4."],
    "opportunities":["기회 문장1.", "기회 문장2.", "기회 문장3.", "기회 문장4."],
    "threats":      ["위협 문장1.", "위협 문장2.", "위협 문장3.", "위협 문장4."]
  }},
  "시장현황": [
    "시장 분석 문장1.",
    "시장 분석 문장2.",
    "시장 분석 문장3.",
    "시장 분석 문장4.",
    "시장 분석 문장5."
  ],
  "핵심경쟁력": [
    {{"제목": "포인트1 키워드", "내용": ["분석 문장1.", "분석 문장2."]}},
    {{"제목": "포인트2 키워드", "내용": ["분석 문장1.", "분석 문장2."]}},
    {{"제목": "포인트3 키워드", "내용": ["분석 문장1.", "분석 문장2."]}}
  ],
  "자금계획": [
    {{"항목": "세부항목명1", "금액": "OOO만원", "내용": ["사용처 문장1.", "사용처 문장2."]}},
    {{"항목": "세부항목명2", "금액": "OOO만원", "내용": ["사용처 문장1.", "사용처 문장2."]}}
  ],
  "매출전망": [
    {{"단계": "1단계 (도입)", "목표": "OOO만원", "내용": ["전략 문장1."]}},
    {{"단계": "2단계 (성장)", "목표": "OOO만원", "내용": ["전략 문장1."]}},
    {{"단계": "3단계 (확장)", "목표": "OOO만원", "내용": ["전략 문장1."]}},
    {{"단계": "4단계 (안착)", "목표": "OOO만원", "내용": ["전략 문장1."]}}
  ],
  "성장비전": {{
    "단기": ["단기 비전 문장1.", "단기 비전 문장2."],
    "중기": ["중기 비전 문장1.", "중기 비전 문장2."],
    "장기": ["장기 비전 문장1.", "장기 비전 문장2."]
  }},
  "조언": [
    "인증 조언 문장1.",
    "인증 조언 문장2.",
    "지식재산권 전략 문장1.",
    "지식재산권 전략 문장2."
  ]
}}"""

    return html_prompt, json_prompt


# ==========================================
# 메인 앱
# ==========================================
def show_main_app():

    with st.sidebar:
        st.markdown("### 🤖 AI 컨설팅 시스템")
        st.markdown("---")
        st.header("⚙️ AI 엔진 설정")

        if "api_key" not in st.session_state:
            st.session_state["api_key"] = st.secrets.get("GEMINI_API_KEY", "")

        if st.session_state["api_key"]:
            st.success("✅ API KEY 저장됨")
            genai.configure(api_key=st.session_state["api_key"])
            if st.button("🔄 API KEY 변경", use_container_width=True):
                st.session_state["api_key"] = ""
                st.rerun()
        else:
            api_input = st.text_input("Gemini API Key", type="password",
                                      placeholder="API 키를 입력하세요")
            if st.button("💾 API KEY 저장", use_container_width=True):
                if api_input:
                    st.session_state["api_key"] = api_input
                    genai.configure(api_key=api_input)
                    st.success("✅ 저장됨!")
                    time.sleep(0.5)
                    st.rerun()
                else:
                    st.warning("API 키를 입력해주세요.")

        st.markdown("---")
        st.header("📂 업체 관리")
        db = load_db()

        if st.button("💾 현재 정보 저장", use_container_width=True):
            c_name = st.session_state.get("in_company_name","").strip()
            if c_name:
                current = {k:v for k,v in st.session_state.items() if k.startswith("in_")}
                db[c_name] = current
                save_db(db)
                st.success(f"✅ '{c_name}' 저장 완료!")
            else:
                st.warning("기업명을 먼저 입력해주세요.")

        selected = st.selectbox("저장된 업체 목록", ["선택 안 함"]+list(db.keys()))
        cs1, cs2 = st.columns(2)
        with cs1:
            if st.button("📂 불러오기", use_container_width=True):
                if selected != "선택 안 함":
                    for k, v in db[selected].items():
                        st.session_state[k] = v
                    st.rerun()
        with cs2:
            if st.button("🔄 초기화", use_container_width=True):
                for k in list(st.session_state.keys()):
                    if k.startswith("in_"): del st.session_state[k]
                st.rerun()

        st.markdown("---")
        st.header("🚀 빠른 리포트 생성")
        if st.button("📊 1. 기업분석리포트 생성", use_container_width=True):
            st.session_state["permanent_data"] = {k:v for k,v in st.session_state.items() if k.startswith("in_")}
            st.session_state["view_mode"] = "REPORT"; st.rerun()
        if st.button("💡 2. 정책자금 매칭 리포트", use_container_width=True):
            st.session_state["permanent_data"] = {k:v for k,v in st.session_state.items() if k.startswith("in_")}
            st.session_state["view_mode"] = "MATCHING"; st.rerun()
        if st.button("📝 3. 사업계획서 생성", use_container_width=True):
            st.session_state["permanent_data"] = {k:v for k,v in st.session_state.items() if k.startswith("in_")}
            st.session_state["view_mode"] = "PLAN"; st.rerun()

    if "view_mode" not in st.session_state: st.session_state["view_mode"] = "INPUT"
    if "permanent_data" not in st.session_state: st.session_state["permanent_data"] = {}

    # ─────────────────────────────────────────
    # REPORT 모드
    # ─────────────────────────────────────────
    if st.session_state["view_mode"] == "REPORT":
        if st.button("⬅️ 대시보드로 돌아가기"):
            for k, v in st.session_state["permanent_data"].items():
                st.session_state[k] = v
            st.session_state["view_mode"] = "INPUT"; st.rerun()

        d = st.session_state["permanent_data"]
        c_name = d.get('in_company_name','미입력').strip()
        st.title("📋 시각화 기반 AI 기업분석 리포트")
        st.subheader(f"📌 분석 대상 기업: {c_name}")

        if not st.session_state.get("api_key",""):
            st.error("⚠️ 좌측 사이드바에 API 키를 입력해주세요.")
            return

        # 캐시 재사용
        if "report_response_text" in st.session_state:
            response_text  = st.session_state["report_response_text"]
            fig            = st.session_state["report_fig"]
            monthly_vals   = st.session_state.get("report_monthly_vals", [])
            monthly_labels = st.session_state.get("report_monthly_labels", [])
            if "[GRAPH_INSERT_POINT]" in response_text:
                parts = response_text.partition("[GRAPH_INSERT_POINT]")
                st.markdown(parts[0], unsafe_allow_html=True)
                st.plotly_chart(fig, use_container_width=True)
                st.markdown(parts[2], unsafe_allow_html=True)
            else:
                st.markdown(response_text, unsafe_allow_html=True)
                st.plotly_chart(fig, use_container_width=True)
        else:
            try:
                with st.status("🚀 AI가 리포트 + PPT 데이터를 생성 중입니다...",
                               expanded=True) as status:
                    try:
                        available = [m.name for m in genai.list_models()
                                     if 'generateContent' in m.supported_generation_methods]
                    except Exception as e:
                        raise Exception(f"API 키 권한 오류: {e}")

                    if 'models/gemini-1.5-flash' in available: tm = 'gemini-1.5-flash'
                    elif 'models/gemini-1.5-pro' in available: tm = 'gemini-1.5-pro'
                    elif 'models/gemini-pro' in available: tm = 'gemini-pro'
                    elif available: tm = available[0].replace('models/','')
                    else: raise Exception("사용 가능한 모델이 없습니다.")

                    model = genai.GenerativeModel(tm)
                    html_prompt, json_prompt = generate_report_and_ppt_data(model, c_name, d)

                    # HTML 리포트 생성
                    st.write("📄 HTML 리포트 생성 중...")
                    html_response = model.generate_content(html_prompt)
                    response_text = html_response.text

                    # 후처리: 마침표 줄바꿈 + 하이픈
                    response_text = postprocess_html(response_text)

                    # PPT JSON 생성
                    st.write("📊 PPT 데이터 생성 중...")
                    json_response = model.generate_content(json_prompt)
                    ppt_json_text = json_response.text

                    # JSON 파싱
                    try:
                        clean_json = re.sub(r'```json|```', '', ppt_json_text).strip()
                        ppt_data = json.loads(clean_json)
                    except:
                        ppt_data = {}
                        st.warning("PPT 데이터 파싱 실패 - 기본값으로 생성됩니다.")

                    # 차트 데이터
                    val_cur = safe_int(d.get('in_sales_current',0))
                    if val_cur <= 0: val_cur = 1000
                    sv, ev = val_cur/12, val_cur/12*1.5
                    monthly_vals = []
                    for i in range(12):
                        p = i/11.0
                        monthly_vals.append(
                            int(sv+(ev-sv)*p+(ev-sv)*0.15*np.sin(p*np.pi*3.5)))
                    monthly_labels = [f"{i}월" for i in range(1,13)]

                    fig = go.Figure()
                    fig.add_trace(go.Scatter(
                        x=monthly_labels, y=monthly_vals,
                        mode='lines+markers+text',
                        text=[format_kr_currency(v) for v in monthly_vals],
                        textposition="top center", textfont=dict(size=11),
                        line=dict(color='#0EA5E9', width=4, shape='spline'),
                        marker=dict(size=10, color='#10B981',
                                    line=dict(width=2, color='white'))
                    ))
                    fig.update_layout(
                        title="📈 향후 1년간 월별 예상 매출 상승 곡선",
                        xaxis_title="진행 월", yaxis_title="예상 매출액",
                        xaxis=dict(showgrid=False),
                        yaxis=dict(showgrid=True, gridcolor='#BAE6FD'),
                        template="plotly_white",
                        margin=dict(l=20,r=20,t=40,b=20),
                        paper_bgcolor="rgba(0,0,0,0)",
                        plot_bgcolor="rgba(0,0,0,0)"
                    )

                    status.update(label="✅ 리포트 + PPT 데이터 생성 완료!",
                                  state="complete")

                # 세션 저장
                st.session_state["report_response_text"]  = response_text
                st.session_state["report_fig"]            = fig
                st.session_state["report_d"]              = d
                st.session_state["report_monthly_vals"]   = monthly_vals
                st.session_state["report_monthly_labels"] = monthly_labels
                st.session_state["report_ppt_data"]       = ppt_data

                if "[GRAPH_INSERT_POINT]" in response_text:
                    parts = response_text.partition("[GRAPH_INSERT_POINT]")
                    st.markdown(parts[0], unsafe_allow_html=True)
                    st.plotly_chart(fig, use_container_width=True)
                    st.markdown(parts[2], unsafe_allow_html=True)
                else:
                    st.markdown(response_text, unsafe_allow_html=True)
                    st.plotly_chart(fig, use_container_width=True)

                st.balloons()

            except Exception as e:
                st.error(f"❌ 오류 발생: {str(e)}")
                return

        # ── 다운로드 ──
        st.divider()
        st.subheader("💾 리포트 다운로드")
        response_text  = st.session_state.get("report_response_text","")
        report_d       = st.session_state.get("report_d", d)
        monthly_vals   = st.session_state.get("report_monthly_vals",[])
        monthly_labels = st.session_state.get("report_monthly_labels",[])
        ppt_data       = st.session_state.get("report_ppt_data",{})
        safe_fn = "".join([c for c in c_name
                           if c.isalnum() or c in (" ","_")]).strip() or "업체"

        dl1, dl2 = st.columns(2)
        with dl1:
            st.markdown("**📄 PDF 다운로드**")
            st.caption("HTML 받은 후 → 브라우저 Ctrl+P → PDF 저장 (그래프 포함)")
            html_data = generate_html_download(
                c_name, response_text, report_d,
                monthly_vals, monthly_labels)
            st.download_button(
                label="📥 HTML 다운로드 (그래프 포함)",
                data=html_data,
                file_name=f"{safe_fn}_기업분석.html",
                mime="text/html; charset=utf-8",
                type="primary",
                use_container_width=True
            )
        with dl2:
            st.markdown("**📊 PPT 다운로드**")
            st.caption("하늘+민트 클린 모던 / 10슬라이드 / 내용 100% 반영")
            if st.button("📊 PPT 생성하기", use_container_width=True):
                with st.spinner("PPT 생성 중... (10슬라이드)"):
                    pptx_data = generate_pptx(
                        c_name, ppt_data, report_d,
                        monthly_vals, monthly_labels)
                if pptx_data:
                    st.download_button(
                        label="📥 PPT 다운로드",
                        data=pptx_data,
                        file_name=f"{safe_fn}_기업분석.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        type="primary",
                        use_container_width=True
                    )

    # ─────────────────────────────────────────
    # MATCHING 모드
    # ─────────────────────────────────────────
    elif st.session_state["view_mode"] == "MATCHING":
        if st.button("⬅️ 대시보드로 돌아가기"):
            for k, v in st.session_state["permanent_data"].items():
                st.session_state[k] = v
            st.session_state["view_mode"] = "INPUT"; st.rerun()

        d = st.session_state["permanent_data"]
        c_name = d.get('in_company_name','미입력').strip()
        st.title("🎯 AI 정책자금 최적화 매칭 리포트")
        st.subheader(f"📌 분석 대상 기업: {c_name}")

        if not st.session_state.get("api_key",""):
            st.error("⚠️ 좌측 사이드바에 API 키를 입력해주세요.")
            return

        try:
            with st.status("🚀 AI가 심사를 진행 중입니다...", expanded=True) as status:
                available = [m.name for m in genai.list_models()
                             if 'generateContent' in m.supported_generation_methods]
                tm = 'gemini-1.5-flash' if 'models/gemini-1.5-flash' in available else 'gemini-pro'
                model = genai.GenerativeModel(tm)

                total_debt = format_kr_currency(sum([safe_int(d.get(k,0)) for k in
                    ['in_debt_kosme','in_debt_semas','in_debt_koreg','in_debt_kodit',
                     'in_debt_kibo','in_debt_etc','in_debt_credit','in_debt_coll']]))
                s_25 = format_kr_currency(safe_int(d.get('in_sales_2025',0)))
                c_ind = d.get('in_industry','미입력')
                biz_type = d.get('in_biz_type','개인')
                nice_score = safe_int(d.get('in_nice_score',0))
                req_fund = format_kr_currency(safe_int(d.get('in_req_amount',0)))
                cert = "보유" if any([d.get('in_chk_6'),d.get('in_chk_4'),
                                       d.get('in_chk_10')]) else "미보유"
                biz_years = 0
                if d.get('in_start_date','').strip():
                    try: biz_years = max(0,2026-int(d.get('in_start_date','')[:4]))
                    except: pass

                prompt = f"""당신은 전문 경영컨설턴트입니다.
규칙: 마크다운 금지. '~있음','~예상됨' 간결체. '~습니다' 금지. <br> 금지.
[입력] 기업명:{c_name}/업종:{c_ind}/전년도매출:{s_25}/총기대출:{total_debt}/필요자금:{req_fund}

<h2 style="color:#0369A1;border-bottom:2px solid #0369A1;padding-bottom:8px;margin-top:30px;">1. 기업 스펙 진단 요약</h2>
<div style="background:#EFF6FF;padding:18px;border-radius:13px;border:1px solid #BAE6FD;margin-bottom:13px;">
<b>기업명:</b> {c_name} | <b>업종:</b> {c_ind} ({biz_type}) | <b>업력:</b> 약 {biz_years}년 | <b>NICE:</b> {nice_score}점 | <b>인증:</b> {cert}<br>
<b>전년도매출:</b> <span style="color:#0369A1;font-weight:bold;">{s_25}</span> | <b>총 기대출:</b> <span style="color:#EF4444;">{total_debt}</span> | <b>필요자금: {req_fund}</b>
</div>
<div style="background:#ECFDF5;padding:13px;border-radius:10px;border-left:4px solid #10B981;margin-bottom:17px;line-height:2.0;">(진단 요약 2~3문장)</div>

<h2 style="color:#0369A1;border-bottom:2px solid #0369A1;padding-bottom:8px;margin-top:30px;">2. 우선순위 추천 정책자금</h2>
<table style="width:100%;table-layout:fixed;border-collapse:separate;border-spacing:13px;margin-bottom:13px;">
<tr>
<td style="background:#ECFDF5;padding:18px;border-radius:13px;border-left:4px solid #10B981;vertical-align:top;line-height:2.0;"><b style="color:#10B981;">🥇 1순위: [기관명] / [자금명] / 예상한도</b><div style="margin-top:9px;">(사유 2~3문장)</div></td>
<td style="background:#ECFDF5;padding:18px;border-radius:13px;border-left:4px solid #10B981;vertical-align:top;line-height:2.0;"><b style="color:#10B981;">🥈 2순위: [기관명] / [자금명] / 예상한도</b><div style="margin-top:9px;">(사유 2~3문장)</div></td>
</tr>
</table>

<h2 style="color:#0369A1;border-bottom:2px solid #0369A1;padding-bottom:8px;margin-top:30px;">3. 후순위 추천</h2>
<table style="width:100%;table-layout:fixed;border-collapse:separate;border-spacing:13px;margin-bottom:13px;">
<tr>
<td style="background:#FFFBEB;padding:18px;border-radius:13px;border-left:4px solid #F59E0B;vertical-align:top;line-height:2.0;"><b style="color:#F59E0B;">🥉 3순위: [기관명] / [자금명] / 예상한도</b><div style="margin-top:9px;">(사유 2~3문장)</div></td>
<td style="background:#FFFBEB;padding:18px;border-radius:13px;border-left:4px solid #F59E0B;vertical-align:top;line-height:2.0;"><b style="color:#F59E0B;">🏅 4순위: [기관명] / [자금명] / 예상한도</b><div style="margin-top:9px;">(사유 2~3문장)</div></td>
</tr>
</table>

<h2 style="color:#0369A1;border-bottom:2px solid #0369A1;padding-bottom:8px;margin-top:30px;">4. 보완 가이드</h2>
<div style="background:#FEF2F2;border-left:4px solid #EF4444;padding:18px;border-radius:13px;line-height:2.0;">
<b style="color:#EF4444;">🚨 보완 조언</b>
<div style="margin-top:9px;">(전략 2~3문장)</div>
</div>"""
                response = model.generate_content(prompt)
                status.update(label="✅ 매칭 리포트 생성 완료!", state="complete")

            processed = postprocess_html(response.text)
            st.markdown(processed, unsafe_allow_html=True)
            st.balloons()
            st.divider()
            safe_fn = "".join([c for c in c_name
                               if c.isalnum() or c in (" ","_")]).strip() or "업체"
            html_bytes = f"""<!DOCTYPE html><html lang="ko"><head><meta charset="UTF-8">
<title>{c_name} 매칭리포트</title>
<style>*{{box-sizing:border-box;-webkit-print-color-adjust:exact!important;print-color-adjust:exact!important;}}
body{{font-family:'Malgun Gothic','Apple SD Gothic Neo',sans-serif;padding:35px;line-height:2.0;
color:#1e293b;max-width:1050px;margin:0 auto;background:#fff;font-size:14px;}}
.print-btn{{display:block;width:100%;padding:13px;
background:linear-gradient(135deg,#0EA5E9,#10B981);
color:white;font-size:16px;font-weight:bold;border:none;
border-radius:10px;cursor:pointer;margin-bottom:26px;text-align:center;}}
@media print{{.print-btn{{display:none;}}@page{{size:A4;margin:12mm;}}
body{{padding:0!important;font-size:12px!important;}}}}</style>
</head><body>
<button class="print-btn" onclick="window.print()">🖨️ 클릭하여 PDF로 저장</button>
<h1 style="color:#0369A1;text-align:center;">🎯 AI 정책자금 매칭 리포트: {c_name}</h1>
{processed}</body></html>""".encode('utf-8')
            st.download_button(
                label="📥 매칭 리포트 다운로드 (HTML→PDF)",
                data=html_bytes,
                file_name=f"{safe_fn}_매칭리포트.html",
                mime="text/html; charset=utf-8",
                type="primary"
            )
        except Exception as e:
            st.error(f"❌ 오류 발생: {str(e)}")

    # ─────────────────────────────────────────
    # PLAN 모드
    # ─────────────────────────────────────────
    elif st.session_state["view_mode"] == "PLAN":
        if st.button("⬅️ 대시보드로 돌아가기"):
            for k, v in st.session_state["permanent_data"].items():
                st.session_state[k] = v
            st.session_state["view_mode"] = "INPUT"; st.rerun()

        d = st.session_state["permanent_data"]
        c_name     = d.get('in_company_name','미입력').strip()
        rep_name   = d.get('in_rep_name','미입력')
        c_ind      = d.get('in_industry','미입력')
        career     = d.get('in_career','미입력')
        s_25       = format_kr_currency(d.get('in_sales_2025',0))
        total_debt = format_kr_currency(sum([safe_int(d.get(k,0)) for k in
            ['in_debt_kosme','in_debt_semas','in_debt_koreg','in_debt_kodit',
             'in_debt_kibo','in_debt_etc','in_debt_credit','in_debt_coll']]))
        nice       = d.get('in_nice_score',0)
        item       = d.get('in_item_desc','미입력')
        market     = d.get('in_market_status','미입력')
        diff       = d.get('in_diff_point','미입력')
        cert       = "보유" if any([d.get('in_chk_6'),d.get('in_chk_4'),
                                     d.get('in_chk_10')]) else "미보유"
        req_fund   = format_kr_currency(d.get('in_req_amount',0))
        fund_type  = d.get('in_fund_type','운전자금')
        fund_purpose = d.get('in_fund_purpose','미입력')
        biz_years  = 0
        if d.get('in_start_date','').strip():
            try: biz_years = max(0,2026-int(d.get('in_start_date','')[:4]))
            except: pass

        ds = (f"[기업] {c_name}/{rep_name}/{c_ind}/업력{biz_years}년/{career}\n"
              f"[재무] 전년매출:{s_25}/기대출:{total_debt}/NICE:{nice}점\n"
              f"[비즈니스] {item}/{market}/{diff}/인증:{cert}\n"
              f"[자금] {req_fund}({fund_type}/{fund_purpose})")
        gem = "https://gemini.google.com/app/"
        prompts = {
            "kosme_plan": f"중진공 심사역. 사업계획서 초안. 포커스: 고용창출, 기술성, 미래성장성.\n{ds}",
            "kosme_loan": f"중진공 심사역. 융자신청서 초안. 포커스: 재무분석, 자금조달 및 상환계획.\n{ds}",
            "semas_plan": f"소진공 심사역. 사업계획서 초안. 포커스: 사업생존가능성, 지역상권 영업전략.\n{ds}",
            "semas_loan": f"소진공 심사역. 융자신청서 초안. 포커스: 매출대비 고정비 및 상환능력 증빙.\n{ds}",
            "kodit_plan": f"신보 심사역. 사업계획서 초안. 포커스: 매출 J커브 성장세, 차별성→매출확대.\n{ds}",
            "kodit_loan": f"신보 심사역. 보증신청서 초안. 포커스: 기대출 대비 유동성 해결 및 상환능력.\n{ds}",
            "kibo_plan":  f"기보 심사역. 기술사업계획서 초안. 포커스: 기술혁신성, 특허현황, R&D역량.\n{ds}",
            "kibo_loan":  f"기보 심사역. 기술평가 보증신청서 초안. 포커스: 기술개발자금 사용처, 상용화 후 재무성과.\n{ds}",
            "ir_plan":    f"VC 심사역. PSST 사업계획서 초안. 포커스: Problem, Solution, Scale-up, Team.\n{ds}",
            "ir_loan":    f"VC 심사역. 투자제안 1-Pager 초안. 포커스: 핵심가치, 자금필요성, Exit시나리오.\n{ds}",
        }
        st.title("📝 기관별 맞춤형 Gems 프롬프트 팩")
        st.info("아래 프롬프트를 복사하여 각 기관별 Gems 주소로 들어가 붙여넣기 하세요.")
        tabs = st.tabs(["1. 중진공","2. 소진공","3. 신보/재단","4. 기술보증기금","5. 제안용(IR)"])
        cfgs = [
            ("kosme_plan","kosme_loan","🏢 중소벤처기업진흥공단"),
            ("semas_plan","semas_loan","🏪 소상공인시장진흥공단"),
            ("kodit_plan","kodit_loan","🏦 신용보증기금 / 지역신보"),
            ("kibo_plan", "kibo_loan", "🔬 기술보증기금"),
            ("ir_plan",   "ir_loan",   "📈 제안용 (IR / PSST)"),
        ]
        for tab, (k1,k2,title) in zip(tabs, cfgs):
            with tab:
                st.subheader(title)
                c1, c2 = st.columns(2)
                with c1:
                    st.link_button("🚀 사업계획서 Gems", gem, use_container_width=True)
                    st.code(prompts[k1], language="markdown")
                with c2:
                    st.link_button("📝 융자/보증신청서 Gems", gem, use_container_width=True)
                    st.code(prompts[k2], language="markdown")

    # ─────────────────────────────────────────
    # INPUT 모드 (대시보드)
    # ─────────────────────────────────────────
    else:
        for key in ["report_response_text","report_fig","report_d",
                    "report_monthly_vals","report_monthly_labels","report_ppt_data"]:
            if key in st.session_state: del st.session_state[key]

        st.title("📊 AI 컨설팅 대시보드")
        c1,c2,c3 = st.columns(3)
        with c1:
            if st.button("📊 1. 기업분석리포트 생성", use_container_width=True, type="primary"):
                st.session_state["permanent_data"] = {k:v for k,v in st.session_state.items() if k.startswith("in_")}
                st.session_state["view_mode"] = "REPORT"; st.rerun()
        with c2:
            if st.button("💡 2. 정책자금 매칭 리포트", use_container_width=True, type="primary"):
                st.session_state["permanent_data"] = {k:v for k,v in st.session_state.items() if k.startswith("in_")}
                st.session_state["view_mode"] = "MATCHING"; st.rerun()
        with c3:
            if st.button("📝 3. 사업계획서 생성", use_container_width=True, type="primary"):
                st.session_state["permanent_data"] = {k:v for k,v in st.session_state.items() if k.startswith("in_")}
                st.session_state["view_mode"] = "PLAN"; st.rerun()

        st.markdown("<br>", unsafe_allow_html=True)
        st.header("1. 기업현황")
        c1,c2,c3 = st.columns(3)
        with c1:
            st.text_input("기업명", key="in_company_name")
            st.text_input("사업자번호", key="in_raw_biz_no")
            bt = st.radio("사업자유형", ["개인","법인"], horizontal=True, key="in_biz_type")
            if bt == "법인": st.text_input("법인등록번호", key="in_raw_corp_no")
        with c2:
            st.text_input("사업개시일", placeholder="2020.01.01", key="in_start_date")
            st.selectbox("업종", ["제조업","서비스업","IT업","도소매업","건설업","기타"], key="in_industry")
            ls = st.radio("사업장 임대여부", ["자가","임대"], horizontal=True, key="in_lease_status")
            if ls == "임대":
                lc1,lc2 = st.columns(2)
                with lc1: st.number_input("보증금(만원)", value=0, step=1, key="in_lease_deposit")
                with lc2: st.number_input("월임대료(만원)", value=0, step=1, key="in_lease_rent")
        with c3:
            st.text_input("전화번호", key="in_biz_tel")
            st.text_input("사업장 주소", key="in_biz_addr")
            ha = st.radio("추가사업장현황", ["무","유"], horizontal=True, key="in_has_additional_biz")
            if ha == "유": st.text_input("추가 사업장 정보", key="in_additional_biz_addr")

        st.markdown("<br>", unsafe_allow_html=True)
        st.header("2. 대표자 정보")
        r1,r2,r3 = st.columns(3)
        with r1:
            rc1,rc2 = st.columns(2)
            with rc1: st.text_input("대표자명", key="in_rep_name")
            with rc2: st.text_input("생년월일", key="in_rep_dob")
            st.text_input("연락처", key="in_rep_phone")
            st.selectbox("통신사", ["SKT","KT","LG U+","알뜰폰"], key="in_rep_telecom")
            st.text_input("이메일 주소", key="in_rep_email")
        with r2:
            st.text_input("거주지 주소", key="in_home_addr")
            st.radio("거주지 상태", ["자가","임대"], horizontal=True, key="in_home_status")
            st.text_input("부동산 현황", key="in_real_estate")
        with r3:
            st.text_input("최종학교", key="in_edu_school")
            st.text_input("학과", key="in_edu_major")
            st.text_area("경력(최근기준)", key="in_career")

        st.markdown("<br>", unsafe_allow_html=True)
        st.header("3. 신용 및 연체 정보")
        cr1,cr2 = st.columns(2)
        with cr1:
            cc1,cc2 = st.columns(2)
            with cc1: st.radio("세금체납", ["무","유"], horizontal=True, key="in_tax_status")
            with cc2: st.radio("금융연체", ["무","유"], horizontal=True, key="in_fin_status")
            sc1,sc2 = st.columns(2)
            with sc1: kcb = st.number_input("KCB 점수", value=0, step=1, key="in_kcb_score")
            with sc2: nice = st.number_input("NICE 점수", value=0, step=1, key="in_nice_score")
        with cr2:
            st.info(f"#### 🏆 등급 판정 결과\n\n"
                    f"* **KCB:** {get_credit_grade(kcb,'KCB')}등급\n"
                    f"* **NICE:** {get_credit_grade(nice,'NICE')}등급")

        st.markdown("<br>", unsafe_allow_html=True)
        st.header("4. 재무현황")
        m1,m2,m3,m4 = st.columns(4)
        with m1: st.number_input("금년 매출(만원)", value=0, step=1, key="in_sales_current")
        with m2: st.number_input("25년도 매출합계(만원)", value=0, step=1, key="in_sales_2025")
        with m3: st.number_input("24년도 매출합계(만원)", value=0, step=1, key="in_sales_2024")
        with m4: st.number_input("23년도 매출합계(만원)", value=0, step=1, key="in_sales_2023")

        st.markdown("<br>", unsafe_allow_html=True)
        st.header("5. 기대출현황")
        d1,d2,d3,d4 = st.columns(4)
        with d1: st.number_input("중진공(만원)", value=0, step=1, key="in_debt_kosme")
        with d2: st.number_input("소진공(만원)", value=0, step=1, key="in_debt_semas")
        with d3: st.number_input("신용보증재단(만원)", value=0, step=1, key="in_debt_koreg")
        with d4: st.number_input("신용보증기금(만원)", value=0, step=1, key="in_debt_kodit")
        d5,d6,d7,d8 = st.columns(4)
        with d5: st.number_input("기술보증기금(만원)", value=0, step=1, key="in_debt_kibo")
        with d6: st.number_input("신용대출(만원)", value=0, step=1, key="in_debt_credit")
        with d7: st.number_input("담보대출(만원)", value=0, step=1, key="in_debt_coll")
        with d8: st.number_input("기타(만원)", value=0, step=1, key="in_debt_etc")

        st.markdown("<br>", unsafe_allow_html=True)
        st.header("6. 필요자금")
        p1,p2,p3 = st.columns([1,1,2])
        with p1: st.selectbox("자금구분", ["운전자금","시설자금"], key="in_fund_type")
        with p2: st.number_input("필요자금액(만원)", value=0, step=1, key="in_req_amount")
        with p3: st.text_input("자금사용용도", key="in_fund_purpose")

        st.markdown("<br>", unsafe_allow_html=True)
        st.header("7. 인증현황")
        ac1,ac2,ac3,ac4 = st.columns(4)
        with ac1: st.checkbox("소상공인확인서", key="in_chk_1"); st.checkbox("창업확인서", key="in_chk_2")
        with ac2: st.checkbox("여성기업확인서", key="in_chk_3"); st.checkbox("이노비즈", key="in_chk_4")
        with ac3: st.checkbox("벤처인증", key="in_chk_6"); st.checkbox("뿌리기업확인서", key="in_chk_7")
        with ac4: st.checkbox("ISO인증", key="in_chk_10")

        st.markdown("<br>", unsafe_allow_html=True)
        st.header("8. 비즈니스 정보")
        st.text_area("[아이템]", key="in_item_desc")
        st.markdown("**[주거래처 정보]**")
        cli1,cli2,cli3 = st.columns(3)
        with cli1: st.text_input("거래처 1", key="in_client_1")
        with cli2: st.text_input("거래처 2", key="in_client_2")
        with cli3: st.text_input("거래처 3", key="in_client_3")
        st.text_area("[판매루트]", key="in_sales_route")
        st.text_area("[시장현황]", key="in_market_status")
        st.text_area("[차별화]", key="in_diff_point")
        st.text_area("[앞으로의 계획]", key="in_future_plan")
        st.markdown("<br>", unsafe_allow_html=True)
        st.success("✅ 세팅 완료! 좌측에 API 키 저장하시고 상단 버튼을 클릭해 주십시오.")

show_main_app()
